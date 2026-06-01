#!/usr/bin/env python3
"""
Genera la plantilla editable del PPT del Proyecto 4.

Uso:
    python scripts/generate_ppt.py --grupo 4

Crea: presentacion/Proyecto 4 - G4.pptx
"""
import argparse
import os
import sys

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
sys.path.insert(0, ROOT)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Instala python-pptx: pip install python-pptx")
    sys.exit(1)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0


def build_ppt(grupo: int, estudiantes: str) -> str:
    os.makedirs(os.path.join(ROOT, "presentacion"), exist_ok=True)
    out_path = os.path.join(ROOT, "presentacion", f"Proyecto 4 - G{grupo}.pptx")

    prs = Presentation()
    add_title_slide(
        prs,
        "Proyecto 4 — Clasificación con IA",
        f"Grupo {grupo}\n{estudiantes}\nPontificia Universidad Javeriana",
    )

    slides_content = [
        (
            "Problema",
            [
                "¿Predecir intento de suicidio vs ideación suicida en Bogotá?",
                "Clasificación binaria con relevancia clínica (screening).",
                "Variables: sociodemográficas y factores de riesgo psicosociales.",
            ],
        ),
        (
            "Base de datos",
            [
                "Fuente: Datos Abiertos Bogotá — Secretaría Distrital de Salud.",
                "Repositorio público tipo benchmark (datos reales de salud).",
                "~254.000 registros | Target: clasificaciondelaconducta.",
                "Insertar: tabla de variables y enlace al dataset.",
            ],
        ),
        (
            "EDA y limpieza",
            [
                "Análisis exploratorio (balance, nulos, correlaciones).",
                "Verificación: duplicados, target válido, outliers en edad.",
                "Limpieza: estandarización educativa, imputación, export CSV.",
                "Insertar: capturas de notebooks/01_eda_limpieza.ipynb.",
            ],
        ),
        (
            "Diseño experimental",
            [
                "Factorial 2³ por técnica → 8 configs × 3 técnicas = 24 runs.",
                "RF: n_estimators, max_depth, min_samples_split.",
                "SVM: C, kernel, gamma (muestra 15.000 por costo O(n²)).",
                "MLP: hidden_layer_sizes, activation, learning_rate_init.",
                "Insertar: tabla desde results/experiment_table.csv.",
            ],
        ),
        (
            "Entrenamiento — Random Forest",
            [
                "sklearn.ensemble.RandomForestClassifier",
                "Insertar: captura de entrenamiento / métricas del notebook 03.",
            ],
        ),
        (
            "Entrenamiento — SVM",
            [
                "sklearn.svm.SVC con kernel RBF/poly",
                "Insertar: captura del notebook 03.",
            ],
        ),
        (
            "Entrenamiento — MLP",
            [
                "sklearn.neural_network.MLPClassifier",
                "Insertar: curva de pérdida y captura de entrenamiento.",
            ],
        ),
        (
            "Matrices de confusión (mejor modelo por técnica)",
            [
                "Selección por F1 ponderado (desbalance 73/27).",
                "Insertar: results/confusion_matrices/cm_combinadas_mejores.png",
            ],
        ),
        (
            "Tabla y análisis comparativo",
            [
                "Insertar: results/experiment_table.csv",
                "Insertar: resumen de results/analisis_comparativo.md",
                "Comparar accuracy, F1 y sensibilidad a hiperparámetros.",
            ],
        ),
        (
            "Conclusiones",
            [
                "Mejor técnica según métrica elegida (accuracy / F1).",
                "Limitaciones: SVM con submuestra; interpretación clínica.",
                "Trabajo futuro: más features temporales, validación externa.",
            ],
        ),
        (
            "Bibliografía",
            [
                "Scikit-learn: Pedregosa et al., JMLR 2011.",
                "Datos Abiertos Bogotá — conducta suicida.",
                "Material del curso Introducción a la IA — diseño experimental.",
            ],
        ),
    ]

    for title, bullets in slides_content:
        add_bullet_slide(prs, title, bullets)

    prs.save(out_path)
    return out_path


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--grupo", type=int, required=True, help="Número de grupo (X en Proyecto 4 - GX)")
    parser.add_argument(
        "--estudiantes",
        default="Nombre1 Apellido1\nNombre2 Apellido2",
        help="Nombres para portada y comentario de entrega",
    )
    args = parser.parse_args()
    path = build_ppt(args.grupo, args.estudiantes)
    print(f"✅ PPT creado: {path}")
    print(f"Comentario buzón:\nGrupo {args.grupo}\n{args.estudiantes}")


if __name__ == "__main__":
    main()
